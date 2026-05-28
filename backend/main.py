import json
import io
import os
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
import anthropic
import fitz
from pptx import Presentation

app = FastAPI(title="DealLens API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

SNAPSHOT_SCHEMA = """Return ONLY valid JSON, no markdown fences, no extra text:
{
  "companyName": "string",
  "oneLiner": "1 sentence: product + customer + problem solved",
  "stage": "e.g. Seed / Series A",
  "sector": "e.g. FinTech / AI",
  "location": "City, State",
  "verdict": "Pass" or "Explore" or "Fast-track",
  "overview": "3-4 sentence company overview",
  "team": [{"name": "string", "role": "string", "background": "1-2 sentences"}],
  "market": {
    "tam": "e.g. $65B",
    "sam": "e.g. $12B",
    "som": "e.g. $800M",
    "trends": "2-3 sentences on market dynamics and why now"
  },
  "traction": [{"label": "string", "value": "string", "note": "optional context"}],
  "competitors": [{"name": "string", "positioning": "string", "differentiation": "how startup differs"}],
  "financials": {
    "raising": "e.g. $25M",
    "valuation": "e.g. $180M pre",
    "round": "e.g. Series A",
    "useOfFunds": ["e.g. 40% — Engineering team expansion"]
  },
  "redFlags": ["specific concern"],
  "questions": ["specific question to ask founder"],
  "verdictReasoning": "3-4 sentence verdict with key reasoning"
}"""

MEMO_SECTIONS = [
    "Executive Summary",
    "Company & Product Overview",
    "Market Opportunity & Timing",
    "Competitive Landscape",
    "Team Assessment",
    "Financial Model & Use of Funds",
    "Risk Analysis",
    "Recommendation & Next Steps",
]


def extract_pdf(data: bytes) -> str:
    doc = fitz.open(stream=data, filetype="pdf")
    return "\n".join(page.get_text() for page in doc).strip()


def extract_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    lines = []
    for i, slide in enumerate(prs.slides):
        lines.append(f"--- Slide {i + 1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


@app.post("/analyze")
async def analyze(file: UploadFile = File(...)):
    raw = await file.read()
    name = (file.filename or "").lower()

    if name.endswith(".pdf"):
        text = extract_pdf(raw)
    elif name.endswith((".pptx", ".ppt")):
        text = extract_pptx(raw)
    else:
        raise HTTPException(400, "Upload a PDF or PPTX file")

    if not text.strip():
        raise HTTPException(400, "No text could be extracted from this file")

    text = text[:50000]

    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4096,
        system="You are a senior VC analyst. Analyze pitch decks and return structured JSON deal snapshots. Be direct and opinionated. If data is missing, infer from context or note it. Return valid JSON only — no markdown, no explanation.",
        messages=[{
            "role": "user",
            "content": f"Analyze this pitch deck and return a deal snapshot as JSON.\n\nRequired schema:\n{SNAPSHOT_SCHEMA}\n\nPitch deck content:\n{text}",
        }],
    )

    raw_text = msg.content[0].text.strip()
    # Strip markdown code fences if Claude adds them
    if raw_text.startswith("```"):
        parts = raw_text.split("```")
        raw_text = parts[1].lstrip("json").strip() if len(parts) > 1 else raw_text

    try:
        return json.loads(raw_text)
    except json.JSONDecodeError:
        raise HTTPException(500, "Failed to parse analysis — try again")


@app.post("/generate-memo")
async def generate_memo(snapshot: dict):
    sections = "\n".join(f"{i+1}. {s}" for i, s in enumerate(MEMO_SECTIONS))
    msg = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=8000,
        system="You are a VC analyst writing a partner-ready investment memo. Be structured, direct, and opinionated. Format output as HTML using <h2> for section headings, <p> for body paragraphs, <ul><li> for lists. Do NOT include <html>, <body>, or <head> tags.",
        messages=[{
            "role": "user",
            "content": f"Write a full investment memo with these sections as <h2> headings:\n{sections}\n\nBase it on this deal snapshot:\n{json.dumps(snapshot, indent=2)}",
        }],
    )
    return {"html": msg.content[0].text}


# Serve frontend from repo root
frontend = Path(__file__).parent.parent
app.mount("/", StaticFiles(directory=str(frontend), html=True), name="static")
